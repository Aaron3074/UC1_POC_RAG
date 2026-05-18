# =============================================================================
# ENTERPRISE KNOWLEDGE ASSISTANT — RAG Engine v7
# Pure Python · rank-bm25 · python-docx · openpyxl · python-pptx
# No API keys · No LLM · No cloud dependencies
#
# v7 fixes over v6 (targeted, no architectural changes):
#
#   FIX 1 — CHUNK DEDUPLICATION (Priority 1)
#     Root cause: sub-query splitting introduced in v6 scored some chunks
#     under multiple sub-queries and used chunk_idx as the dedup key. When
#     chunk_idx was None (common in window-chunked docs), the fallback was
#     the loop index `i` — which changes per sub-query run. Same chunk got
#     different dedup keys per pass and was inserted multiple times.
#     Fix: stable composite key (doc_id, content_fingerprint). Content
#     fingerprint is hash(chunk_content[:120]) — invariant across runs.
#
#   FIX 2 — CAPACITY BREACH → ANALYTICAL, NOT INCIDENT
#     Root cause: bare `breach` pattern in INCIDENT fired before ANALYTICAL
#     for queries like "which system is closest to a capacity breach".
#     Fix: (a) Narrow INCIDENT's breach pattern to explicit security/SLA
#     breach phrases only. (b) Add capacity-specific patterns to ANALYTICAL
#     (closest to, approaching, near limit). (c) Move ANALYTICAL before
#     INCIDENT in priority_order for capacity-class queries via a two-pass
#     detect: if ANALYTICAL fires on the first pass, return it immediately
#     without reaching INCIDENT.
#
#   FIX 3 — LIST SYNTHESIS FALLBACK FOR RISK/APPROVAL QUERIES
#     Root cause: _list() only pattern-matched on structured prefixes
#     (INC-XXXX, REL-XXX, etc.). Open-ended risk queries have no such
#     prefixes, so _list() fell through to _summary(), missing substantive
#     risk content. Fix: detect risk/approval query vocabulary; when
#     matched, fall through to _analytical() instead of _summary().
#     When structured items AND risk content both exist, render both.
#
#   FIX 4 — PARAMETER EXTRACTOR NOISE (WEEKEND, SYMPTOMS, etc.)
#     Root cause: param_pat matched any ALLCAPS token followed by : or =,
#     even with whitespace between token and delimiter. Informal docs with
#     section headers like "CONTACTS THIS WEEKEND:" matched because WEEKEND
#     was followed by ":". Fix: (a) Tighten regex so colon must be followed
#     immediately by a non-whitespace character (actual value, not newline).
#     (b) Add a blocklist of common English words that appear all-caps in
#     informal handover/ops docs.
#
#   ADDITIONAL: 8 new unstructured knowledge base documents added covering
#     DBA notes, fraud model constraints, network VLAN layout, platform retro,
#     security notes, new-analyst FAQ, data-eng/ETL notes, SWIFT compliance.
#
# Architecture is unchanged: same 7-phase pipeline, same class hierarchy.
# =============================================================================

import os
import re
import json
import math
import logging
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Tuple, Set
from collections import defaultdict, Counter

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger('KnowledgeAssistant')


# =============================================================================
# AUTO SYNONYM DISCOVERY ENGINE
# =============================================================================

class AutoSynonymDiscovery:
    """
    Discovers synonym/alias relationships from raw document text using three signals:

    Signal 1 — ACRONYM EXPANSION
      Detects "Full Name (ABBREV)" and "ABBREV (Full Name)" patterns.
      e.g. "Core Banking System (CBS)" → {cbs: [core banking system]}
           "the SSO (Single Sign-On) service" → {sso: [single sign-on]}

    Signal 2 — CO-OCCURRENCE CLUSTERING
      Finds pairs of terms that appear in the same sentence repeatedly.
      If "payments gateway" and "payments-gw" co-occur in 3+ sentences,
      they are likely synonyms. Threshold configurable.

    Signal 3 — FILENAME NORMALISATION
      Derives term aliases from filenames/folder names.
      "payment_gateway_spec.docx" → ["payment gateway", "payment_gateway_spec"]

    Results are merged with the manual registry (manual always wins on conflicts).
    Exported to auto_synonyms_discovered.yaml for analyst review.
    """

    # Matches "Full Multi Word Name (ABBR)" or "ABBR (Full Name)"
    ACRONYM_PATTERN = re.compile(
        r'\b([A-Z][a-zA-Z\s\-]{3,40})\s*\(([A-Z][A-Z0-9\-]{1,15})\)'  # Full (ABBR)
        r'|'
        r'\b([A-Z][A-Z0-9\-]{1,15})\s*\(([A-Za-z][a-zA-Z\s\-]{3,40})\)'  # ABBR (full)
    )

    SENT_SPLIT = re.compile(r'(?<=[.!?])\s+|\n')

    def __init__(self, co_occur_threshold: int = 2):
        self.co_occur_threshold = co_occur_threshold
        self.discovered: Dict[str, Set[str]] = defaultdict(set)

    def discover_from_documents(self, documents: dict) -> Dict[str, List[str]]:
        """Run all discovery signals across all documents. Returns {canonical: [alias, ...]}."""
        logger.info("=" * 65)
        logger.info("AUTO-SYNONYM DISCOVERY")
        logger.info("=" * 65)

        all_sentences = []
        all_filenames = []

        for doc in documents.values():
            text = doc.get("raw_text", "")
            self._extract_acronyms(text)
            sents = self.SENT_SPLIT.split(text)
            all_sentences.extend([s.strip() for s in sents if len(s.strip()) > 20])
            all_filenames.append((doc.get("filename", ""), doc.get("topic_folder", "")))

        self._extract_cooccurrence(all_sentences)

        for fname, folder in all_filenames:
            self._extract_from_filename(fname, folder)

        result = {}
        for canonical, aliases in self.discovered.items():
            clean = [a for a in aliases if a != canonical and len(a) > 1]
            if clean:
                result[canonical] = sorted(clean)

        logger.info(f"  Discovered {len(result)} synonym groups from documents")
        for k, v in list(result.items())[:8]:
            logger.info(f"    {k} → {v[:4]}")
        if len(result) > 8:
            logger.info(f"    ... and {len(result)-8} more")
        logger.info("=" * 65)
        return result

    def _normalise(self, text: str) -> str:
        return re.sub(r'[\s\-_]+', ' ', text.strip().lower())

    def _extract_acronyms(self, text: str):
        for match in self.ACRONYM_PATTERN.finditer(text):
            if match.group(1) and match.group(2):
                full = self._normalise(match.group(1))
                abbr = match.group(2).strip().lower()
                if len(abbr) >= 2 and len(full) > len(abbr):
                    self.discovered[abbr].add(full)
                    self.discovered[full].add(abbr)
            elif match.group(3) and match.group(4):
                abbr = match.group(3).strip().lower()
                full = self._normalise(match.group(4))
                if len(abbr) >= 2 and len(full) > len(abbr):
                    self.discovered[abbr].add(full)
                    self.discovered[full].add(abbr)

    def _extract_cooccurrence(self, sentences: List[str]):
        """
        Count how often capitalised/acronym-like token pairs co-occur in the same sentence.
        Pairs above co_occur_threshold are candidates for synonym relationships.
        """
        token_pat = re.compile(
            r'\b([A-Z][A-Z0-9\-]{1,15}|[A-Z][a-z]+(?:[\s\-][A-Z][a-z]+){1,4})\b'
        )
        pair_counts: Counter = Counter()
        for sent in sentences:
            tokens = list(set(t.lower() for t in token_pat.findall(sent)))
            tokens = [t for t in tokens if len(t) > 2 and t not in {
                'the', 'for', 'and', 'this', 'that', 'with', 'from', 'are',
                'all', 'not', 'can', 'has', 'have', 'was', 'been', 'will',
            }]
            for i, a in enumerate(tokens):
                for b in tokens[i+1:]:
                    if a != b:
                        pair = tuple(sorted([a, b]))
                        pair_counts[pair] += 1

        for (a, b), count in pair_counts.items():
            if count >= self.co_occur_threshold:
                # Only link tokens of comparable length (avoids spurious cross-domain pairs)
                if abs(len(a) - len(b)) < max(len(a), len(b)) * 0.7:
                    self.discovered[a].add(b)
                    self.discovered[b].add(a)

    def _extract_from_filename(self, filename: str, folder: str):
        stem  = re.sub(r'\.(txt|docx|xlsx|pptx)$', '', filename.lower())
        parts = re.split(r'[_\-\s]+', stem)
        noise = {'the', 'a', 'an', 'and', 'or', 'of', 'to', 'for', 'in',
                 'with', 'by', 'from', 'on', 'at', 'informal', 'formal',
                 'guide', 'notes', 'log', 'doc', 'spec', 'register', 'report',
                 'v1', 'v2', 'v3', 'v4', 'v5', 'v6', 'v7', 'final', 'draft',
                 'new', 'old'}
        meaningful = [p for p in parts if len(p) > 2 and p not in noise]
        if len(meaningful) >= 2:
            joined = ' '.join(meaningful)
            for part in meaningful:
                if len(part) > 3:
                    self.discovered[joined].add(part)

    def export_yaml(self, path: Path, discovered: Dict[str, List[str]]):
        """
        Write discovered synonyms to YAML for analyst review.
        Analysts promote useful entries to synonym_registry.yaml.
        """
        lines = [
            "# AUTO-DISCOVERED SYNONYMS",
            f"# Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            "# Review these and promote useful ones to synonym_registry.yaml",
            "# Lines starting with # are comments.",
            "",
        ]
        for canonical, aliases in sorted(discovered.items()):
            if len(canonical) > 2 and aliases:
                lines.append(f"{canonical}:")
                for alias in aliases[:6]:
                    lines.append(f"  - {alias}")
        path.write_text("\n".join(lines), encoding="utf-8")
        logger.info(f"  Auto-discovered synonyms exported to: {path}")


# =============================================================================
# SYNONYM REGISTRY  (manual + auto-discovered, merged)
# =============================================================================

class SynonymRegistry:
    """
    Loads manual synonym_registry.yaml and merges auto-discovered synonyms.
    Manual entries always win on conflicts — auto-discovery never overwrites
    a manually curated entry.
    """

    def __init__(self, registry_path: str = None):
        self.canonical_to_aliases: Dict[str, List[str]] = {}
        self.alias_to_canonical:   Dict[str, str]       = {}
        self._manual_canonicals: Set[str] = set()

        if registry_path is None:
            registry_path = Path(__file__).parent / "synonym_registry.yaml"
        self._registry_path = Path(registry_path)
        self._load_manual(self._registry_path)

    def _load_manual(self, path: Path):
        if not path.exists():
            logger.warning(f"  Manual synonym registry not found: {path}")
            return
        current_key = None
        try:
            with open(path, encoding="utf-8") as f:
                for line in f:
                    stripped = line.rstrip()
                    if not stripped or stripped.lstrip().startswith("#"):
                        continue
                    if not line.startswith(" ") and not line.startswith("\t") and ":" in line:
                        current_key = line.split(":")[0].strip().lower()
                        self.canonical_to_aliases[current_key] = []
                        self._manual_canonicals.add(current_key)
                    elif stripped.lstrip().startswith("- ") and current_key:
                        alias = stripped.lstrip()[2:].strip().lower()
                        if alias:
                            self.canonical_to_aliases[current_key].append(alias)
                            self.alias_to_canonical[alias] = current_key
        except Exception as e:
            logger.warning(f"  Failed to load manual registry: {e}")

        logger.info(f"  Manual synonym registry: {len(self.canonical_to_aliases)} groups, "
                    f"{len(self.alias_to_canonical)} aliases")

    def merge_auto_discovered(self, discovered: Dict[str, List[str]]):
        """
        Merge auto-discovered synonyms. Manual entries always win on conflict.
        New groups are added; existing manual groups are only extended with
        aliases not already present.
        """
        added_groups  = 0
        added_aliases = 0

        for canonical, aliases in discovered.items():
            if canonical in self._manual_canonicals:
                continue
            if canonical in self.alias_to_canonical:
                continue
            if len(canonical) < 3:
                continue

            if canonical not in self.canonical_to_aliases:
                clean_aliases = [
                    a for a in aliases
                    if a != canonical and len(a) > 2
                    and a not in self.alias_to_canonical
                ]
                if clean_aliases:
                    self.canonical_to_aliases[canonical] = clean_aliases
                    for alias in clean_aliases:
                        self.alias_to_canonical[alias] = canonical
                    added_groups  += 1
                    added_aliases += len(clean_aliases)
            else:
                existing = set(self.canonical_to_aliases[canonical])
                for alias in aliases:
                    if alias not in existing and alias != canonical and len(alias) > 2:
                        if alias not in self.alias_to_canonical:
                            self.canonical_to_aliases[canonical].append(alias)
                            self.alias_to_canonical[alias] = canonical
                            existing.add(alias)
                            added_aliases += 1

        logger.info(f"  Auto-synonyms merged: +{added_groups} new groups, +{added_aliases} aliases")
        logger.info(f"  Total registry: {len(self.canonical_to_aliases)} groups, "
                    f"{len(self.alias_to_canonical)} aliases")

    def expand(self, text: str) -> str:
        """Return text with all recognised synonyms appended."""
        text_lower = text.lower()
        expansions = []
        for canonical, aliases in self.canonical_to_aliases.items():
            all_terms = [canonical] + aliases
            if any(term in text_lower for term in all_terms):
                for term in all_terms:
                    if term not in text_lower:
                        expansions.append(term)
        return (text + " " + " ".join(expansions)) if expansions else text

    def expand_tokens(self, tokens: List[str]) -> List[str]:
        """Token-level expansion for BM25 indexing."""
        extra = []
        for tok in tokens:
            if tok in self.canonical_to_aliases:
                extra.extend(self.canonical_to_aliases[tok])
            elif tok in self.alias_to_canonical:
                canonical = self.alias_to_canonical[tok]
                extra.append(canonical)
                extra.extend(self.canonical_to_aliases.get(canonical, []))
        return tokens + [t for t in extra if t not in tokens]


# =============================================================================
# PHASE 1 — DOCUMENT INGESTION  (resilient, multi-format)
# =============================================================================

class DocumentIngestionEngine:

    SUPPORTED = {".txt", ".docx", ".xlsx", ".pptx"}

    def __init__(self, kb_dir: str, synonyms: SynonymRegistry):
        self.kb_dir   = Path(kb_dir)
        self.synonyms = synonyms
        self.documents: Dict[str, dict] = {}

        logger.info("=" * 65)
        logger.info("PHASE 1: Ingestion Engine v7")
        logger.info(f"  Base: {self.kb_dir}")
        logger.info("=" * 65)

        if not self.kb_dir.exists():
            raise FileNotFoundError(f"Knowledge base not found: {self.kb_dir}")

        self._ingest_all()
        self._print_summary()

    def _ingest_all(self):
        files = []
        for ext in self.SUPPORTED:
            files.extend(self.kb_dir.rglob(f"*{ext}"))
        files = sorted(set(files))
        by_ext = Counter(f.suffix for f in files)
        logger.info(f"  Found {len(files)} files: " +
                    ", ".join(f"{e}:{n}" for e, n in sorted(by_ext.items())))
        for path in files:
            self._load(path)

    def _load(self, path: Path):
        ext = path.suffix.lower()
        try:
            raw = self._parse(path, ext)
            if not raw or len(raw.strip()) < 30:
                logger.warning(f"  SKIP (empty): {path.name}")
                return
            expanded = self.synonyms.expand(raw)
            doc_id   = re.sub(r'[^a-z0-9_]', '_', path.stem.lower()) + ext.replace(".", "_")
            folder   = str(path.relative_to(self.kb_dir).parent)
            meta     = self._extract_metadata(raw, path, folder)
            chunks   = self._chunk(raw, meta)

            self.documents[doc_id] = {
                "doc_id":        doc_id,
                "filename":      path.name,
                "file_format":   ext,
                "topic_folder":  folder,
                "full_path":     str(path),
                "raw_text":      raw,
                "expanded_text": expanded,
                "chunks":        chunks,
                "char_count":    len(raw),
                "metadata":      meta,
            }
            quality = meta.get("_quality", "?")
            logger.info(f"  [{ext:5s}] [{folder}] {path.name} — "
                        f"{len(raw):,}ch  {len(chunks)} chunks  quality:{quality}")
        except Exception as e:
            logger.error(f"  FAILED [{path.name}]: {e}")

    # ── Format parsers ────────────────────────────────────────────────────────

    def _parse(self, path: Path, ext: str) -> str:
        if ext == ".txt":
            return path.read_text(encoding="utf-8", errors="replace")
        elif ext == ".docx":
            return self._parse_docx(path)
        elif ext == ".xlsx":
            return self._parse_xlsx(path)
        elif ext == ".pptx":
            return self._parse_pptx(path)
        return ""

    def _parse_docx(self, path: Path) -> str:
        from docx import Document
        doc   = Document(path)
        lines = []
        for para in doc.paragraphs:
            t = para.text.strip()
            if not t:
                continue
            if para.style.name.startswith("Heading"):
                lvl   = re.search(r'\d+', para.style.name)
                depth = int(lvl.group()) if lvl else 2
                lines.append(f"\n{'=' * max(3, 6-depth)} {t} {'=' * max(3, 6-depth)}")
            else:
                lines.append(t)
        for table in doc.tables:
            lines.append("")
            for row in table.rows:
                cells     = [c.text.strip() for c in row.cells]
                non_empty = [c for c in cells if c and c != "None"]
                if non_empty:
                    lines.append(" | ".join(non_empty))
        return "\n".join(lines)

    def _parse_xlsx(self, path: Path) -> str:
        import openpyxl
        wb    = openpyxl.load_workbook(path, data_only=True)
        lines = []
        for name in wb.sheetnames:
            ws = wb[name]
            lines.append(f"\n=== {name} ===")
            for row in ws.iter_rows(values_only=True):
                cells     = [str(c).strip() if c is not None else "" for c in row]
                non_empty = [c for c in cells if c and c.lower() != "none"]
                if non_empty:
                    lines.append("  " + " | ".join(non_empty))
        return "\n".join(lines)

    def _parse_pptx(self, path: Path) -> str:
        from pptx import Presentation
        prs   = Presentation(path)
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                lines.append(f"\n=== {texts[0]} ===")
                for t in texts[1:]:
                    lines.append(t)
        return "\n".join(lines)

    # ── Metadata extraction ───────────────────────────────────────────────────

    def _extract_metadata(self, text: str, path: Path, folder: str) -> dict:
        meta             = {}
        structured_fields = 0
        for line in text.split("\n")[:15]:
            line = line.strip()
            if re.match(r"^\s*={3,}", line) or line.startswith("|"):
                continue
            if ":" in line:
                parts = line.split(":", 1)
                key   = parts[0].strip().upper().replace(" ", "_")
                value = parts[1].strip()
                if key and value and len(key) < 35 and len(value) < 200:
                    meta[key] = value
                    structured_fields += 1

        meta["_quality"] = (
            "structured"   if structured_fields >= 3 else
            "partial"      if structured_fields >= 1 else
            "unstructured"
        )

        if "DOCUMENT" not in meta:
            stem = re.sub(r'[\-_]+', ' ', path.stem)
            stem = re.sub(r'\b(v\d+|final|FINAL|draft|DRAFT|new|NEW|old|OLD|copy|use_this|informal)\b',
                          '', stem, flags=re.IGNORECASE)
            meta["DOCUMENT"] = stem.strip().title()

        if "TOPIC" not in meta:
            meta["TOPIC"] = folder.replace("_", " ").replace("/", " > ").title()

        if "SOURCE_FORMAT" not in meta:
            meta["SOURCE_FORMAT"] = path.suffix.upper().replace(".", "")

        return meta

    # ── Resilient chunking ────────────────────────────────────────────────────

    CHUNK_SIZE    = 500
    CHUNK_OVERLAP = 100

    def _chunk(self, text: str, meta: dict) -> List[dict]:
        """
        Three-tier fallback: section headers → paragraph boundaries → sliding window.
        Each strategy guarantees at least one chunk with meaningful content.
        """
        chunks = self._try_section_split(text)
        if len(chunks) < 2:
            chunks = self._try_paragraph_split(text)
        if len(chunks) < 2:
            chunks = self._sliding_window(text)
        for i, chunk in enumerate(chunks):
            chunk["chunk_idx"] = i
            chunk["doc_title"] = meta.get("DOCUMENT", "")
            chunk["topic"]     = meta.get("TOPIC", "")
            chunk["owner"]     = meta.get("OWNER", "")
            chunk["quality"]   = meta.get("_quality", "unknown")
        return chunks

    def _try_section_split(self, text: str) -> List[dict]:
        header_pat = re.compile(r'={3,}\s*(.+?)\s*={3,}')
        parts  = header_pat.split(text)
        chunks = []
        if len(parts) > 2:
            i = 1
            while i < len(parts):
                title   = parts[i].strip() if parts[i] else ""
                content = parts[i+1].strip() if (i+1) < len(parts) else ""
                if content and len(content) > 30:
                    chunks.append({"title": title, "content": content, "method": "section"})
                i += 2
        return chunks

    def _try_paragraph_split(self, text: str) -> List[dict]:
        chunks = []
        for para in text.split("\n\n"):
            para = para.strip()
            if len(para) > 80:
                lines = para.split("\n")
                title = lines[0].strip()[:80] if lines else ""
                chunks.append({"title": title, "content": para, "method": "paragraph"})
        return chunks

    def _sliding_window(self, text: str) -> List[dict]:
        chunks = []
        text   = text.strip()
        start  = 0
        while start < len(text):
            end     = min(start + self.CHUNK_SIZE, len(text))
            content = text[start:end].strip()
            if len(content) > 30:
                chunks.append({
                    "title":   f"Section {len(chunks)+1}",
                    "content": content,
                    "method":  "window",
                })
            start += self.CHUNK_SIZE - self.CHUNK_OVERLAP
        return chunks

    def _print_summary(self):
        by_folder  = defaultdict(int)
        by_fmt     = Counter()
        by_quality = Counter()
        for doc in self.documents.values():
            by_folder[doc["topic_folder"]] += 1
            by_fmt[doc["file_format"]]     += 1
            by_quality[doc["metadata"].get("_quality", "?")] += 1

        logger.info("\n  === INGESTION SUMMARY ===")
        logger.info(f"  Documents : {len(self.documents)}")
        logger.info(f"  Formats   : {dict(by_fmt)}")
        logger.info(f"  Quality   : {dict(by_quality)}")
        for f in sorted(by_folder):
            logger.info(f"    [{f}] → {by_folder[f]} file(s)")
        logger.info("  " + "=" * 45)


# =============================================================================
# PHASE 2 — BM25 INDEXING
# =============================================================================

class BM25IndexEngine:

    STOPWORDS = {
        "the","a","an","and","or","but","in","on","at","to","for","of","with",
        "by","from","is","are","was","were","be","been","have","has","had",
        "do","does","did","will","would","could","should","may","might","this",
        "that","these","those","it","its","as","if","not","all","any","each",
        "per","no","only","also","than","then","so","up","out","about","into",
        "through","during","before","after","i","you","he","she","we","they",
        "can","which","who","how","what","when","where","there","must","use",
        "using","well","within","without","every","our","your","their","both",
        "some","such","same","other","more","most","used","been","just","like",
    }

    def __init__(self, documents: dict, synonyms: SynonymRegistry):
        from rank_bm25 import BM25Okapi
        self.documents    = documents
        self.synonyms     = synonyms
        self.chunk_corpus: List[dict] = []
        self._build_corpus()

        chunk_tokens = [c["tokens"] for c in self.chunk_corpus]
        self.bm25_chunks = BM25Okapi(chunk_tokens) if chunk_tokens else None
        doc_tokens   = [self._tokenize(doc["expanded_text"]) for doc in documents.values()]
        self.doc_ids = list(documents.keys())
        self.bm25_docs = BM25Okapi(doc_tokens) if doc_tokens else None

        logger.info("=" * 65)
        logger.info(f"PHASE 2: BM25 Index — {len(self.chunk_corpus)} chunks / {len(documents)} docs")
        logger.info("=" * 65)

    def _tokenize(self, text: str) -> List[str]:
        text   = text.lower()
        text   = re.sub(r"[^a-z0-9\s\-_\./]", " ", text)
        tokens = [t for t in text.split() if len(t) > 1 and t not in self.STOPWORDS]
        return tokens if tokens else ["_empty_"]

    def _build_corpus(self):
        for doc_id, doc in self.documents.items():
            for chunk in doc["chunks"]:
                expanded = self.synonyms.expand(chunk["content"])
                tokens   = self._tokenize(expanded)
                self.chunk_corpus.append({
                    "doc_id":    doc_id,
                    "filename":  doc["filename"],
                    "folder":    doc["topic_folder"],
                    "format":    doc["file_format"],
                    "metadata":  doc["metadata"],
                    "full_path": doc["full_path"],
                    "title":     chunk.get("title", ""),
                    "content":   chunk["content"],
                    "chunk_idx": chunk.get("chunk_idx"),
                    "method":    chunk.get("method", ""),
                    "tokens":    tokens,
                })


# =============================================================================
# PHASE 3 — RETRIEVAL
# =============================================================================

class RetrievalEngine:
    """
    BM25 retrieval with intent-aware folder boosting and sub-query splitting
    for long/complex queries. v7 fixes chunk deduplication using a stable
    composite key so that sub-query merging cannot insert the same chunk twice.
    """

    # Folder score multipliers per intent — inject domain expertise into retrieval ranking
    INTENT_FOLDER_BOOST = {
        "TROUBLESHOOT": {"troubleshooting": 2.0, "unstructured_samples": 1.1},
        "INCIDENT":     {"incident_reports": 2.0, "troubleshooting": 1.3, "unstructured_samples": 1.1},
        "LOOKUP":       {"connectivity_ports": 2.5, "application_parameters": 1.5, "troubleshooting": 1.2},
        "HOWTO":        {"troubleshooting": 1.5, "process_workflows": 1.2},
        "DEPENDENCY":   {"connectivity_ports": 2.0, "network_architecture": 1.3},
        "ANALYTICAL":   {"capacity_planning": 1.5, "sla_contracts": 1.5, "incident_reports": 1.3},
    }

    def __init__(self, index: BM25IndexEngine, synonyms: SynonymRegistry):
        self.index    = index
        self.synonyms = synonyms
        logger.info("=" * 65)
        logger.info("PHASE 3: BM25 Retrieval Engine v7 — Ready")
        logger.info("=" * 65)

    def search(self, query: str, top_k: int = 6,
               folder_filter: Optional[str] = None,
               intent: Optional[str] = None) -> List[dict]:
        """
        Search across all chunks with synonym expansion and intent-aware boosting.

        For queries longer than 12 words, the query is also split into two
        halves and scored independently. Results are merged by taking the
        maximum effective score per document across all sub-queries.

        v7 FIX: chunk deduplication now uses a stable composite key:
            (chunk_idx, content_fingerprint)
        where content_fingerprint = hash(content[:120]). This is invariant
        across sub-query runs, preventing the same chunk from appearing in
        top_chunks more than once per document — eliminating repeated content
        blocks in synthesis output.
        """
        if not self.index.bm25_chunks or not self.index.chunk_corpus:
            return []

        # Build sub-query list: full query always first, plus halves for long queries
        words       = query.split()
        sub_queries = [query]
        if len(words) > 12:
            mid = len(words) // 2
            sub_queries.append(" ".join(words[:mid]))
            sub_queries.append(" ".join(words[mid:]))

        # Accumulate best score per document and all unique chunks per document
        doc_best_score: Dict[str, float] = defaultdict(float)
        doc_chunks:     Dict[str, List]  = defaultdict(list)

        for sq in sub_queries:
            expanded     = self.synonyms.expand(sq)
            tokens       = self.index._tokenize(expanded)
            if not tokens:
                continue
            chunk_scores = self.index.bm25_chunks.get_scores(tokens)

            for i, score in enumerate(chunk_scores):
                if score <= 0:
                    continue
                chunk  = self.index.chunk_corpus[i]
                doc_id = chunk["doc_id"]
                folder = chunk["folder"]

                if folder_filter and folder_filter.lower() not in folder.lower():
                    continue

                # Apply intent folder boost (take highest applicable multiplier)
                boost = 1.0
                if intent and intent in self.INTENT_FOLDER_BOOST:
                    for boosted_folder, factor in self.INTENT_FOLDER_BOOST[intent].items():
                        if boosted_folder in folder.lower():
                            boost = max(boost, factor)

                effective_score = score * boost
                doc_best_score[doc_id] = max(doc_best_score[doc_id], effective_score)

                # ── v7 FIX: stable dedup key ────────────────────────────────
                # chunk_idx alone is unreliable: it can be None for some chunking
                # strategies, and even when set, the fallback `i` (corpus index)
                # differs between sub-query runs, causing duplicates to slip through.
                # The content fingerprint is invariant — same chunk = same hash.
                content_fp = hash(chunk.get("content", "")[:120])
                chunk_idx  = chunk.get("chunk_idx")
                dedup_key  = (chunk_idx, content_fp)

                existing_keys = {
                    (c.get("chunk_idx"), hash(c.get("content", "")[:120]))
                    for c in doc_chunks[doc_id]
                }
                if dedup_key not in existing_keys:
                    doc_chunks[doc_id].append({**chunk, "bm25_score": round(effective_score, 4)})
                else:
                    # Update the score on the existing entry if we found a better one
                    for existing_chunk in doc_chunks[doc_id]:
                        if hash(existing_chunk.get("content", "")[:120]) == content_fp:
                            if effective_score > existing_chunk.get("bm25_score", 0):
                                existing_chunk["bm25_score"] = round(effective_score, 4)
                            break
                # ── end v7 FIX ───────────────────────────────────────────────

        if not doc_best_score:
            return []

        ranked  = sorted(doc_best_score, key=lambda x: doc_best_score[x], reverse=True)
        results = []
        for doc_id in ranked[:top_k]:
            doc    = self.index.documents[doc_id]
            chunks = sorted(doc_chunks[doc_id], key=lambda x: x["bm25_score"], reverse=True)
            results.append({
                **doc,
                "bm25_score":  round(doc_best_score[doc_id], 4),
                "final_score": round(doc_best_score[doc_id], 4),
                "top_chunks":  chunks[:5],   # top 5 unique chunks per document
                "chunk_count": len(chunks),
            })
        return results


# =============================================================================
# PHASE 4 — INTENT DETECTION
# =============================================================================

class IntentDetector:
    """
    Nine intent classes detected via priority-ordered regex patterns.

    v7 changes:
    - INCIDENT: bare `breach` pattern removed; replaced with explicit phrases
      (security breach, data breach, SLA breach). This prevents capacity-related
      queries from routing to INCIDENT when they mention "capacity breach".
    - ANALYTICAL: new patterns for capacity comparison queries (closest to,
      approaching, near limit, which system is most/closest).
    - Priority order: ANALYTICAL is checked before INCIDENT in a first-pass
      scan for capacity-class queries, so "which system is closest to a
      capacity breach" routes correctly to ANALYTICAL.
    """

    PATTERNS = {
        "LOOKUP": [
            r"\bwhat (port|ip|address|url|host|version|value|default|parameter|param|timeout|limit|key|endpoint|threshold|ttl)\b",
            r"\bwhat is the (port|ip|address|url|version|value|default|timeout|limit|threshold)\b",
            r"\bwhere is\b",
            r"\bwhich (port|host|ip|address)\b",
            r"\btell me the\b",
        ],
        "HOWTO": [
            r"\bhow (do i|do we|can i|can we|to)\b",
            r"\bsteps (to|for)\b",
            r"\bprocedure (to|for)\b",
            r"\bprocess (to|for)\b",
        ],
        "COMPARISON": [
            r"\b(vs|versus|compared to|difference between|compare)\b",
            r"\b(pros and cons|trade.?off|similarities|differences)\b",
            r"\bwhy .+ (chosen|selected|picked|preferred)\b",
            r"\bbetter (than|or)\b",
        ],
        "SUMMARY": [
            r"\b(summarise|summarize|summary|overview|explain|describe|tell me about)\b",
            r"\bgive me (a|an) (summary|overview|brief|rundown|breakdown)\b",
            r"\bexplain\b",
            r"\bwhat (is|are|does) .+ (do|mean|cover|include)\b",
        ],
        "LIST": [
            r"\b(list|enumerate|show me all|what are all|which are)\b",
            r"\ball (the|of the)?\s*(apps|applications|services|systems|ports|vendors|tools|documents)\b",
            r"\bgive me (a list|all|every)\b",
        ],
        "TROUBLESHOOT": [
            r"\b(troubleshoot|fix|resolve|debug|diagnose)\b",
            r"\b(not working|failing|failed|broken|down|error|issue|problem|incident|outage)\b",
            r"\b(cannot|can.t|unable to) (log in|login|connect|access|authenticate|reach|process)\b",
            r"\b\d{3} (error|status)\b",
            r"\b(getting|receiving) (an? )?(error|alert|warning|timeout)\b",
        ],
        "INCIDENT": [
            # Explicit incident ID reference always routes to INCIDENT
            r"\binc[-\s]?\d{4}[-\s]?\d{4}\b",
            # Narrative incident review phrases
            r"\b(walk me through|take me through|explain).{0,30}(inc|incident|outage)\b",
            r"\b(what happened|what caused|what triggered|what led to).{0,30}(inc|outage|december)\b",
            r"\b(incident|outage|postmortem|post.mortem|rca|root cause analysis)\b",
            r"\b(timeline|blast radius|impact of|recovery from)\b",
            # ── v7 FIX: bare `breach` removed ──────────────────────────────
            # Previously: r"\b(incident|outage|breach|postmortem|...)\b"
            # `breach` alone matched "capacity breach", routing to INCIDENT
            # instead of ANALYTICAL. Now only explicit security/SLA breach
            # phrases route to INCIDENT.
            r"\b(security breach|data breach|breach of sla|sla breach)\b",
            # ── end v7 FIX ──────────────────────────────────────────────────
        ],
        "ANALYTICAL": [
            r"\b(analyse|analyze|analysis|assess|evaluate)\b",
            r"\b(trend|pattern|contributing factor|why did)\b",
            r"\b(risk|impact|projection|forecast|at risk|concern)\b",
            r"\b(which (system|service|app).{0,20}(most|highest|biggest|worst|best|closest|critical))\b",
            # ── v7 FIX: capacity-specific patterns ─────────────────────────
            # Queries about capacity comparison ("closest to breach/ceiling/limit")
            # now explicitly match ANALYTICAL before INCIDENT can fire.
            r"\b(capacity|utilisation|utilization|headroom|bottleneck|ceiling)\b",
            r"\b(closest to|approaching|near.limit|running out of|near capacity)\b",
            # ── end v7 FIX ──────────────────────────────────────────────────
            r"\b(cost|budget|spend|renewal|contract|vendor)\b",
            r"\b(lessons learned|open items|action items|findings|recommendations)\b",
        ],
        "DEPENDENCY": [
            r"\b(depends on|dependency|dependencies|connects to|integrates with|calls|uses)\b",
            r"\b(upstream|downstream|linked to|flow|chain|integration)\b",
            r"\bwhat (does .+ depend|does .+ connect|does .+ call)\b",
            r"\bwhich (apps|systems|services) (use|call|depend on|connect to)\b",
            r"\b(data flow|request flow|end.to.end|e2e)\b",
        ],
    }

    # Intents that should be tested in first-pass before the standard order.
    # Used to resolve conflicts where two intents share vocabulary.
    # E.g. ANALYTICAL should beat INCIDENT for capacity comparison queries.
    _FIRST_PASS_INTENTS = ["ANALYTICAL"]

    def detect(self, query: str) -> str:
        """
        Two-pass intent detection:

        Pass 1: Check ANALYTICAL first for capacity/risk vocabulary.
          If a strong ANALYTICAL signal fires, return immediately before
          INCIDENT can capture it via its own patterns.

        Pass 2: Standard priority order for all other queries.

        v7 note: This does NOT change routing for genuine incident queries
        (INC-XXXX references, "walk me through the outage") — those patterns
        are specific enough that they only appear in INCIDENT, and Pass 2
        handles them correctly. The two-pass approach only matters when
        both ANALYTICAL and INCIDENT patterns could fire on the same query.
        """
        q = query.lower()

        # Pass 1: ANALYTICAL first-pass for capacity/comparison class
        # Only the capacity/comparison patterns are checked here — not all ANALYTICAL patterns,
        # which would be too broad and could incorrectly capture incident-adjacent queries.
        CAPACITY_PATTERNS = [
            r"\b(capacity|utilisation|utilization|headroom|bottleneck|ceiling)\b",
            r"\b(closest to|approaching|near.limit|running out of|near capacity)\b",
            r"\b(which (system|service|app).{0,20}(most|highest|biggest|worst|best|closest|critical))\b",
        ]
        for pat in CAPACITY_PATTERNS:
            if re.search(pat, q):
                return "ANALYTICAL"

        # Pass 2: Standard priority order
        priority_order = [
            "INCIDENT", "LOOKUP", "TROUBLESHOOT", "DEPENDENCY",
            "COMPARISON", "HOWTO", "LIST", "ANALYTICAL", "SUMMARY",
        ]
        for intent in priority_order:
            for pat in self.PATTERNS.get(intent, []):
                if re.search(pat, q):
                    return intent
        return "SUMMARY"


# =============================================================================
# PHASE 5 — CONFIDENCE SCORING
# =============================================================================

class ConfidenceScorer:
    """
    Five-factor confidence model: document count, top BM25 score, score spread,
    document quality, and exact term match. Returns HIGH / MEDIUM / LOW with
    explicit reasons for each factor — designed for explainability at demos.
    """

    def score(self, query: str, results: List[dict], intent: str) -> dict:
        if not results:
            return {"level": "LOW", "score": 0, "reasons": ["No matching documents found"]}

        reasons = []
        points  = 0
        max_pts = 0

        # Factor 1: number of matching documents (0–25 pts)
        max_pts += 25
        n = len(results)
        if n >= 4:
            points += 25; reasons.append(f"Strong: {n} documents matched the query")
        elif n >= 2:
            points += 15; reasons.append(f"Moderate: {n} documents matched")
        else:
            points += 5;  reasons.append("Weak: only 1 document matched")

        # Factor 2: top BM25 score (0–25 pts)
        max_pts  += 25
        top_score = results[0]["bm25_score"]
        if top_score >= 5.0:
            points += 25; reasons.append(f"Strong relevance score ({top_score:.1f})")
        elif top_score >= 2.0:
            points += 15; reasons.append(f"Moderate relevance score ({top_score:.1f})")
        elif top_score >= 0.5:
            points += 8;  reasons.append(f"Low relevance score ({top_score:.1f}) — partial match")
        else:
            reasons.append(f"Very low relevance score ({top_score:.2f}) — weak match")

        # Factor 3: score spread — is the top source clearly dominant? (0–20 pts)
        max_pts += 20
        if len(results) > 1:
            spread = results[0]["bm25_score"] - results[-1]["bm25_score"]
            if spread >= 3.0:
                points += 20; reasons.append("Top source clearly dominant — high specificity")
            elif spread >= 1.0:
                points += 12; reasons.append("Some score differentiation between sources")
            else:
                points += 5;  reasons.append("Sources scored similarly — topic spans multiple docs")
        else:
            points += 10

        # Factor 4: document quality (0–15 pts)
        max_pts    += 15
        structured  = sum(1 for r in results[:3] if r["metadata"].get("_quality") == "structured")
        if structured >= 2:
            points += 15; reasons.append(f"{structured} of top sources are well-structured documents")
        elif structured >= 1:
            points += 8;  reasons.append("Mixed document quality — some sources are unstructured")
        else:
            reasons.append("Sources are unstructured — review manually for accuracy")

        # Factor 5: exact term match in top source (0–15 pts)
        max_pts     += 15
        q_tokens     = set(query.lower().split())
        top_content  = " ".join(c["content"].lower() for c in results[0].get("top_chunks", [])[:2])
        exact_hits   = sum(1 for t in q_tokens if len(t) > 3 and t in top_content)
        if exact_hits >= 3:
            points += 15; reasons.append(f"Exact query terms found in source ({exact_hits} matches)")
        elif exact_hits >= 1:
            points += 8;  reasons.append(f"Partial query term match ({exact_hits} terms found)")
        else:
            reasons.append("No exact query terms — synonym or semantic match")

        pct   = int((points / max_pts) * 100) if max_pts else 0
        level = "HIGH" if pct >= 70 else "MEDIUM" if pct >= 40 else "LOW"
        return {"level": level, "score": pct, "reasons": reasons}


# =============================================================================
# PHASE 6 — ANSWER SYNTHESIS
# =============================================================================

class AnswerSynthesisEngine:
    """
    Nine intent-specific synthesis strategies, each extracting the most relevant
    content from retrieved chunks and formatting it for operational usefulness.

    v7 changes:
    - _list(): fallback for open-ended risk/approval queries changed from
      _summary() to _analytical(), which does full content extraction rather
      than prefix-pattern matching. When structured items (INC-XXX, REL-XXX)
      AND risk vocabulary both exist, both are rendered.
    - _query_suggestions(): param_pat tightened (colon must be followed
      immediately by a non-whitespace value) and a blocklist of common
      English words that appear ALLCAPS in informal docs (WEEKEND, SYMPTOMS,
      CONTACTS, etc.) is applied before suggestions are generated.
    """

    LINES_PER_CHUNK = 22
    CHUNKS_PER_DOC  = 4
    DOCS_IN_ANSWER  = 5
    CONTEXT_WINDOW  = 6

    def __init__(self):
        self.intent_detector   = IntentDetector()
        self.confidence_scorer = ConfidenceScorer()

    # ── Entry point ───────────────────────────────────────────────────────────

    def synthesise(self, query: str, results: List[dict]) -> dict:
        ts = datetime.now().isoformat()
        if not results:
            return {"query": query, "intent": "UNKNOWN", "status": "NO_RESULTS",
                    "answer": self._no_results(query), "sources": [], "timestamp": ts}

        intent     = self.intent_detector.detect(query)
        confidence = self.confidence_scorer.score(query, results, intent)
        logger.info(f"  Intent: {intent}  |  Confidence: {confidence['level']} ({confidence['score']}%)")

        # Build citation map: filename → citation number (order of appearance)
        cite_map = {}
        for i, r in enumerate(results, 1):
            if r["filename"] not in cite_map:
                cite_map[r["filename"]] = i

        dispatch = {
            "LOOKUP":       self._lookup,
            "HOWTO":        self._howto,
            "COMPARISON":   self._comparison,
            "SUMMARY":      self._summary,
            "LIST":         self._list,
            "TROUBLESHOOT": self._troubleshoot,
            "INCIDENT":     self._incident,
            "ANALYTICAL":   self._analytical,
            "DEPENDENCY":   self._dependency,
        }
        body        = dispatch.get(intent, self._summary)(query, results, cite_map)
        conf_block  = self._confidence_block(confidence)
        suggestions = self._query_suggestions(query, intent, results)
        refs        = self._references(results, cite_map)
        full        = "\n\n".join([body.strip(), conf_block, suggestions, refs])

        return {
            "query":      query,
            "intent":     intent,
            "confidence": confidence,
            "status":     "SUCCESS",
            "answer":     full,
            "sources":    self._source_list(results, cite_map),
            "timestamp":  ts,
        }

    # ── Standard blocks ───────────────────────────────────────────────────────

    def _confidence_block(self, c: dict) -> str:
        bar  = {"HIGH": "██████████", "MEDIUM": "██████░░░░", "LOW": "███░░░░░░░"}
        icon = {"HIGH": "✓", "MEDIUM": "~", "LOW": "!"}
        lvl  = c["level"]
        lines = [
            "─" * 65,
            f"CONFIDENCE: {icon.get(lvl,'?')} {lvl}  {bar.get(lvl,'')}  ({c['score']}%)",
            "─" * 65,
        ]
        for r in c["reasons"]:
            lines.append(f"  • {r}")
        return "\n".join(lines)

    def _query_suggestions(self, query: str, intent: str, results: List[dict]) -> str:
        """
        Extracts system names, incident IDs, parameter names, and release IDs
        from retrieved chunks and builds context-specific follow-up questions.

        v7 FIX — param_pat tightened in two ways:
          1. Colon must be immediately followed by a non-whitespace character.
             This prevents section headers like "CONTACTS THIS WEEKEND:" from
             matching — the colon there is followed by a newline, not a value.
          2. Blocklist: common English words that appear ALL-CAPS in informal
             handover/ops docs (WEEKEND, SYMPTOMS, CONTACTS, etc.) are filtered
             out regardless of whether the pattern matched.
        """
        # Known system identifiers — strict list to avoid noise
        system_pat   = re.compile(
            r'\b(CBS|PAYMENTS-GW|FRAUD-ENGINE|AUTH-SVC|REPORTING-SVC|'
            r'SWIFT-ADAPTER|KAFKA|ETL|DWH|ZABBIX|DATADOG|QRADAR|ELK)\b'
        )
        incident_pat = re.compile(r'\bINC-\d{4}-\d{4}\b')

        # ── v7 FIX: tightened param pattern ────────────────────────────────
        # Old: r'\b([A-Z][A-Z_]{4,30})\s*(?:=|:)'
        #   Problem: \s* allows zero or more whitespace including newlines,
        #   so "CONTACTS THIS WEEKEND:" matched with WEEKEND as the token
        #   (WEEKEND followed by colon with no space between).
        # New: colon variant requires an immediate non-whitespace character after it.
        #   = variant unchanged (assignment always has a value on the same line).
        param_pat = re.compile(
            r'\b([A-Z][A-Z_]{4,30})(?:=|\s*=\s*|:\s*(?=\S))',
            re.MULTILINE
        )
        # Blocklist: words that appear in ALLCAPS in informal docs as section labels,
        # not as actual configuration parameter names.
        PARAM_BLOCKLIST = {
            "WEEKEND", "MONDAY", "TUESDAY", "WEDNESDAY", "THURSDAY",
            "FRIDAY", "SATURDAY", "SUNDAY", "GENERAL", "CONTACTS",
            "SCHEDULED", "SYMPTOMS", "STATUS", "NOTES", "UPDATE",
            "ACTION", "ITEMS", "SUMMARY", "DETAILS", "SECTION",
            "THINGS", "STUFF", "HEADS", "ISSUE", "ISSUES", "STEPS",
            "CHECK", "FIRST", "ESCALATE", "FOLLOW", "LESSON",
        }
        # ── end v7 FIX ─────────────────────────────────────────────────────

        rel_pat = re.compile(r'\bREL-\d{4}-\d{3}\b')
        adr_pat = re.compile(r'\bADR-\d{3}\b')

        systems   = []
        incidents = []
        params    = []
        releases  = []
        adrs      = []

        for r in results[:4]:
            for chunk in r.get("top_chunks", [])[:3]:
                t = chunk["content"]
                systems.extend(system_pat.findall(t))
                incidents.extend(incident_pat.findall(t))
                params.extend(param_pat.findall(t))
                releases.extend(rel_pat.findall(t))
                adrs.extend(adr_pat.findall(t))

        # Deduplicate preserving order; apply blocklist and length filter to params
        systems   = list(dict.fromkeys(systems))[:3]
        incidents = list(dict.fromkeys(incidents))[:2]
        params    = [
            p for p in dict.fromkeys(params)
            if len(p) > 5 and p not in PARAM_BLOCKLIST
        ][:2]
        releases  = list(dict.fromkeys(releases))[:1]
        adrs      = list(dict.fromkeys(adrs))[:1]

        suggestions = []

        if intent == "TROUBLESHOOT":
            if systems:
                suggestions.append(f"What are the SLA targets and availability history for {systems[0]}?")
            if len(systems) > 1:
                suggestions.append(f"Which applications depend on {systems[1]} and what breaks if it goes down?")
            if incidents:
                suggestions.append(f"What action items and fixes came out of {incidents[0]}?")
            elif systems:
                suggestions.append(f"What is the on-call escalation path for a {systems[0]} outage?")

        elif intent == "INCIDENT":
            if incidents:
                suggestions.append(f"What preventive actions were taken after {incidents[0]}?")
                suggestions.append(f"Which other systems were impacted by {incidents[0]}?")
            if len(incidents) > 1:
                suggestions.append(f"How does {incidents[0]} compare in severity to {incidents[1]}?")
            elif systems:
                suggestions.append(f"What is the current capacity status of {systems[0]}?")

        elif intent == "ANALYTICAL":
            if systems:
                suggestions.append(f"What is the full SLA breach history for {systems[0]}?")
            if len(systems) > 1:
                suggestions.append(f"Which upcoming releases affect {systems[1]}?")
            if incidents:
                suggestions.append(f"What patterns caused {incidents[0]} and how were they addressed?")
            elif params:
                suggestions.append(f"What is the impact of changing {params[0]} and who needs to approve it?")

        elif intent == "DEPENDENCY":
            if systems:
                suggestions.append(f"What port and host does {systems[0]} expose its API on?")
            if len(systems) > 1:
                suggestions.append(f"What happens to downstream services if {systems[1]} goes down?")
            if len(systems) > 2:
                suggestions.append(f"How does data flow from {systems[0]} to {systems[2]}?")
            elif systems:
                suggestions.append(f"What are the SLA targets for {systems[0]}?")

        elif intent == "LOOKUP":
            if systems:
                suggestions.append(f"What other configuration parameters does {systems[0]} have?")
            if params:
                suggestions.append(f"What is the impact of changing {params[0]} and does it require CAB approval?")
            if len(systems) > 1:
                suggestions.append(f"How does {systems[0]} connect to {systems[1]}?")
            elif systems:
                suggestions.append(f"What are the SLA targets for {systems[0]}?")

        elif intent == "COMPARISON":
            if adrs:
                suggestions.append(f"What are the consequences documented in {adrs[0]}?")
            if systems:
                suggestions.append(f"What is the current implementation status of {systems[0]}?")
            suggestions.append("What architecture decisions are still open and pending ARB approval?")

        elif intent == "HOWTO":
            if systems:
                suggestions.append(f"Who is the escalation contact for {systems[0]} issues?")
            if incidents:
                suggestions.append(f"What was the root cause of {incidents[0]} and how was it prevented?")
            if params:
                suggestions.append(f"What other parameters interact with {params[0]}?")
            elif systems and len(systems) > 1:
                suggestions.append(f"What are the SLA targets for {systems[1]}?")

        else:
            # Generic content-driven suggestions as fallback
            if incidents:
                suggestions.append(f"What was the full root cause and remediation for {incidents[0]}?")
            if systems:
                suggestions.append(f"What is the current capacity and SLA status of {systems[0]}?")
            if params:
                suggestions.append(f"What happens if {params[0]} is changed without CAB approval?")
            if len(suggestions) < 3 and releases:
                suggestions.append(f"What changes were introduced in {releases[0]} and were there any rollbacks?")
            if len(suggestions) < 3 and len(systems) > 1:
                suggestions.append(f"How do {systems[0]} and {systems[1]} interact with each other?")

        suggestions = suggestions[:3]
        lines = ["─" * 65, "RELATED QUESTIONS", "─" * 65]
        for i, s in enumerate(suggestions, 1):
            lines.append(f"  {i}. {s}")
        return "\n".join(lines)

    def _references(self, results: List[dict], cite_map: dict) -> str:
        """
        Rich reference block: evidence summary, quality warnings,
        matched section titles, and log-normalised relevance percentages.
        """
        lines     = ["─" * 65, "REFERENCES  —  open source files for complete detail", "─" * 65]
        seen      = set()
        max_score = max((r.get("bm25_score", 0) for r in results), default=1)

        for r in results:
            fname = r["filename"]
            if fname in seen:
                continue
            seen.add(fname)
            num   = cite_map.get(fname, "?")
            meta  = r["metadata"]
            qual  = meta.get("_quality", "?")
            raw   = r.get("bm25_score", 0)
            rel   = int(100 * math.log1p(raw) / math.log1p(max_score)) if max_score > 0 else 0
            fmt   = r.get("file_format", "?").upper()
            title = meta.get("DOCUMENT", fname)
            owner = meta.get("OWNER", "—")
            ver   = meta.get("VERSION", "")
            upd   = meta.get("LAST_UPDATED", "")

            # Matched evidence: section titles from top chunks (skip generic "Section N" titles)
            top_chunks = r.get("top_chunks", [])
            evidence   = [
                c.get("title", "") for c in top_chunks[:3]
                if c.get("title") and not c["title"].startswith("Section ")
                and len(c.get("title", "")) > 3
            ]

            lines.append(f"[{num}] {fname}  ({fmt})")
            lines.append(f"     {'─' * 55}")
            lines.append(f"     Title     : {title}")
            lines.append(f"     Folder    : {r['topic_folder']}")
            lines.append(f"     Owner     : {owner}")
            info_parts = []
            if ver: info_parts.append(f"v{ver}")
            if upd: info_parts.append(f"Updated {upd}")
            if info_parts:
                lines.append(f"     Info      : {' | '.join(info_parts)}")
            lines.append(f"     Quality   : {qual}  |  BM25: {raw:.1f}  |  Relevance: {rel}%")
            if evidence:
                lines.append(f"     Matched   : {' › '.join(evidence[:3])}")
            if qual == "unstructured":
                lines.append(f"     ⚠ Note    : Unstructured doc — no formal metadata, verify content manually")
            elif qual == "partial":
                lines.append(f"     ℹ Note    : Partially structured — some metadata inferred from filename/folder")
            lines.append(f"     Path      : {r['full_path']}")
            lines.append("")

        return "\n".join(lines)

    def _source_list(self, results: List[dict], cite_map: dict) -> List[dict]:
        seen, out = set(), []
        max_score = max((r.get("bm25_score", 0) for r in results), default=1)
        for r in results:
            if r["filename"] not in seen:
                seen.add(r["filename"])
                raw = r.get("bm25_score", 0)
                rel = int(100 * math.log1p(raw) / math.log1p(max_score)) if max_score > 0 else 0
                out.append({
                    "citation":      cite_map.get(r["filename"], "?"),
                    "filename":      r["filename"],
                    "file_format":   r.get("file_format"),
                    "topic_folder":  r["topic_folder"],
                    "full_path":     r["full_path"],
                    "bm25_score":    raw,
                    "relevance_pct": rel,
                    "quality":       r["metadata"].get("_quality"),
                    "metadata":      r["metadata"],
                })
        return out

    def _no_results(self, query: str) -> str:
        return (
            f"ANSWER\n{'─' * 65}\n"
            "No relevant documents found for this query.\n\n"
            "Try:\n"
            "  • Rephrasing with different terminology\n"
            "  • Checking the knowledge base covers this topic\n"
            "  • Using folder_filter to narrow scope\n"
            "  • Adding more documents and reindexing"
        )

    # ── Shared helpers ────────────────────────────────────────────────────────

    def _cite(self, fname: str, cite_map: dict) -> str:
        return f"[{cite_map.get(fname, '?')}]"

    def _clean_lines(self, content: str, limit: int = None) -> List[str]:
        lim = limit or self.LINES_PER_CHUNK
        return [
            l.strip() for l in content.split("\n")
            if l.strip() and not re.match(r"^\s*={3,}", l)
        ][:lim]

    def _title(self, r: dict) -> str:
        return r["metadata"].get("DOCUMENT", r["filename"])

    def _extract_matching(self, text: str, keywords: List[str], ctx: int = None) -> List[str]:
        ctx   = ctx or self.CONTEXT_WINDOW
        lines = text.split("\n")[12:]
        seen, out = set(), []
        for i, line in enumerate(lines):
            if re.match(r"^\s*={3,}", line):
                continue
            if any(kw in line.lower() for kw in keywords):
                start = max(0, i - ctx)
                end   = min(len(lines), i + ctx + 1)
                block = "\n".join(l for l in lines[start:end]
                                  if not re.match(r"^\s*={3,}", l)).strip()
                if block and block not in seen:
                    seen.add(block)
                    out.append(block)
        return out

    # ── LOOKUP ────────────────────────────────────────────────────────────────

    def _lookup(self, query: str, results: List[dict], cite_map: dict) -> str:
        """
        Phrase-aware matching for port/parameter/IP queries.
        Detects the lookup type from query vocabulary, re-ranks chunks by how
        well they answer the specific lookup type, then extracts lines that
        contain both the subject AND the value type.
        """
        q_lower = query.lower()
        lines   = [f"ANSWER\n{'─' * 65}"]

        is_port  = any(w in q_lower for w in ["port", "ports"])
        is_param = any(w in q_lower for w in ["parameter", "param", "config",
                                               "setting", "value", "default",
                                               "threshold", "timeout", "ttl"])
        is_host  = any(w in q_lower for w in ["host", "ip", "address", "url", "endpoint"])
        is_ver   = any(w in q_lower for w in ["version", "release"])

        PORT_PAT  = re.compile(r'\b(?:port|port:)\s*\d{2,5}\b|\b\d{4,5}\s*(?:tcp|udp|http|https)\b', re.IGNORECASE)
        PARAM_PAT = re.compile(r'\b[A-Z][A-Z_]{3,}\s*(?:=|:|\|)\s*.+', re.MULTILINE)
        HOST_PAT  = re.compile(r'\b(?:10|172|192)\.\d+\.\d+\.\d+\b|\blocalhost\b|https?://\S+')
        VER_PAT   = re.compile(r'\bv?\d+\.\d+(?:\.\d+)?\b')
        NUM_PAT   = re.compile(r'\b\d+\b')

        stop_lookup = {"port", "ports", "parameter", "param", "config", "host", "ip",
                       "address", "url", "version", "what", "use", "uses", "the", "does",
                       "is", "are", "default", "value", "for", "of", "a", "an", "how"}
        subject_tokens = [w for w in re.findall(r'\w+', q_lower)
                          if w not in stop_lookup and len(w) > 2]

        def chunk_relevance_score(chunk_content: str) -> int:
            """Score a chunk on subject presence + value-type presence."""
            score         = 0
            content_lower = chunk_content.lower()
            score        += sum(2 for t in subject_tokens if t in content_lower)
            if is_port  and PORT_PAT.search(chunk_content):  score += 5
            if is_param and PARAM_PAT.search(chunk_content): score += 5
            if is_host  and HOST_PAT.search(chunk_content):  score += 5
            if is_ver   and VER_PAT.search(chunk_content):   score += 3
            return score

        scored_chunks = []
        for r in results[:5]:
            for chunk in r.get("top_chunks", [])[:4]:
                rel = chunk_relevance_score(chunk["content"])
                if rel > 0:
                    scored_chunks.append((rel, chunk, r))
        scored_chunks.sort(key=lambda x: x[0], reverse=True)

        found = []
        for rel, chunk, r in scored_chunks[:6]:
            cite = self._cite(r["filename"], cite_map)
            target_lines = []
            for line in chunk["content"].split("\n"):
                l = line.strip()
                if not l:
                    continue
                l_lower       = l.lower()
                subject_match = any(t in l_lower for t in subject_tokens)
                value_match   = (
                    (is_port  and PORT_PAT.search(l))  or
                    (is_param and (PARAM_PAT.search(l) or NUM_PAT.search(l))) or
                    (is_host  and HOST_PAT.search(l))  or
                    (is_ver   and VER_PAT.search(l))   or
                    (not any([is_port, is_param, is_host, is_ver]) and subject_match)
                )
                if subject_match and value_match:
                    target_lines.append(l)

            if target_lines:
                block_key = "|".join(target_lines[:3])
                if block_key not in found:
                    found.append(block_key)
                    for tl in target_lines[:6]:
                        lines.append(tl)
                    lines.append(f"  ↳ {cite} {r['filename']}"
                                 + (f" › {chunk['title']}" if chunk.get('title') else ""))
                    lines.append("")
            elif rel > 3 and chunk["content"] not in found:
                found.append(chunk["content"])
                lines.append(f"{cite} {r['filename']}"
                             + (f" › {chunk['title']}" if chunk.get('title') else ""))
                for cl in self._clean_lines(chunk["content"], 12):
                    lines.append(f"  {cl}")
                lines.append("")

        if not lines[1:]:
            for r in results[:2]:
                cite = self._cite(r["filename"], cite_map)
                for chunk in r.get("top_chunks", [])[:2]:
                    lines.append(f"{cite} {r['filename']}"
                                 + (f" › {chunk['title']}" if chunk.get('title') else ""))
                    for cl in self._clean_lines(chunk["content"], 12):
                        lines.append(f"  {cl}")
                    lines.append("")

        if len(results) > 1:
            lines.append("─" * 40)
            lines.append("ADDITIONAL CONTEXT")
            shown = 0
            for r in results[1:4]:
                cite = self._cite(r["filename"], cite_map)
                for chunk in r.get("top_chunks", [])[:1]:
                    if chunk_relevance_score(chunk["content"]) > 1:
                        if chunk.get("title"):
                            lines.append(f"\n{cite} {chunk['title']}  [{r['filename']}]")
                        for cl in self._clean_lines(chunk["content"], 8):
                            lines.append(f"  {cl}")
                        shown += 1
                if shown >= 2:
                    break

        return "\n".join(lines)

    # ── HOWTO ─────────────────────────────────────────────────────────────────

    def _howto(self, query: str, results: List[dict], cite_map: dict) -> str:
        step_pat = re.compile(r'^\s*(STEP\s*\d+[:\-]?|\d+\.\s+|\d+\)\s+)', re.IGNORECASE)
        prereq_pat = re.compile(
            r'^(?:prereq|prerequisite|before you|before running|required:|ensure|dependency:|you need|first[,:])',
            re.IGNORECASE
        )
        note_kws = ["note:","warning:","caution:","important:","do not","never","always"]
        esc_kws  = ["escalat","oncall","pagerduty","@","page","sev-"]

        lines = [f"HOW-TO GUIDE\n{'─' * 65}"]
        prereqs, step_groups, notes, escalation = [], [], [], []

        for r in results[:self.DOCS_IN_ANSWER]:
            cite = self._cite(r["filename"], cite_map)
            for chunk in r["top_chunks"][:self.CHUNKS_PER_DOC]:
                sec_steps = []
                for line in chunk["content"].split("\n"):
                    l = line.strip()
                    if not l:
                        continue
                    l_low = l.lower()
                    if step_pat.match(l):
                        sec_steps.append(l)
                    elif prereq_pat.match(l) and l not in prereqs:
                        prereqs.append(l)
                    elif any(kw in l_low for kw in note_kws) and l not in notes:
                        notes.append(l)
                    elif any(kw in l_low for kw in esc_kws) and l not in escalation:
                        escalation.append(l)
                if sec_steps:
                    step_groups.append({
                        "cite": cite, "source": r["filename"],
                        "section": chunk.get("title", ""), "steps": sec_steps,
                    })

        if prereqs:
            lines += ["", "PREREQUISITES", "─" * 40]
            for p in prereqs[:4]:
                lines.append(f"  • {p}")

        if step_groups:
            lines += ["", "STEPS", "─" * 40]
            for grp in step_groups[:4]:
                if grp["section"]:
                    lines.append(f"\n  {grp['cite']} {grp['section']}  [{grp['source']}]")
                for step in grp["steps"][:20]:
                    lines.append(f"  {step}")
                lines.append("")
        else:
            lines += ["", "GUIDANCE", "─" * 40]
            for r in results[:3]:
                cite = self._cite(r["filename"], cite_map)
                lines.append(f"\n{cite} {self._title(r)}  [{r['topic_folder']}]")
                for chunk in r["top_chunks"][:2]:
                    if chunk.get("title"):
                        lines.append(f"  [{chunk['title']}]")
                    for l in self._clean_lines(chunk["content"]):
                        lines.append(f"    {l}")
                lines.append("")

        if notes:
            lines += ["", "NOTES AND WARNINGS", "─" * 40]
            for n in notes[:5]:
                lines.append(f"  ⚠  {n}")

        if escalation:
            lines += ["", "ESCALATION / CONTACTS", "─" * 40]
            for e in list(dict.fromkeys(escalation))[:5]:
                lines.append(f"  {e}")

        return "\n".join(lines)

    # ── INCIDENT ──────────────────────────────────────────────────────────────

    def _incident(self, query: str, results: List[dict], cite_map: dict) -> str:
        time_kws   = ["timeline","02:","03:","04:","09:","14:","23:","01:","—","→","at "]
        cause_kws  = ["root cause","contributing","caused by","reason","because","due to"]
        impact_kws = ["affected","customer impact","revenue","transactions","users","sla","breach"]
        fix_kws    = ["resolution","fix","restored","resolved","deployed","rolled back","action"]
        action_kws = ["[done]","[open]","action item","prevention","follow.up","lesson"]

        lines = [f"INCIDENT ANALYSIS\n{'─' * 65}"]

        for r in results[:4]:
            cite  = self._cite(r["filename"], cite_map)
            title = self._title(r)
            lines += [f"\n{'─' * 40}", f"{cite} {title}  [{r['topic_folder']}]", ""]

            timeline, causes, impacts, fixes, actions = [], [], [], [], []

            for chunk in r["top_chunks"][:self.CHUNKS_PER_DOC]:
                for line in chunk["content"].split("\n"):
                    l = line.strip()
                    if not l:
                        continue
                    l_low = l.lower()
                    if any(kw in l_low for kw in time_kws) and re.search(r'\d{2}:\d{2}', l):
                        timeline.append(l)
                    elif any(kw in l_low for kw in cause_kws):
                        causes.append(l)
                    elif any(kw in l_low for kw in impact_kws):
                        impacts.append(l)
                    elif any(kw in l_low for kw in fix_kws):
                        fixes.append(l)
                    elif any(kw in l_low for kw in action_kws):
                        actions.append(l)

            for chunk in r["top_chunks"][:2]:
                if chunk.get("title"):
                    lines.append(f"  [{chunk['title']}]")
                for l in self._clean_lines(chunk["content"], 22):
                    lines.append(f"    {l}")
                lines.append("")

            if timeline:
                lines += ["  TIMELINE:", *[f"    {t}" for t in timeline[:15]], ""]
            if causes:
                lines += ["  ROOT CAUSE:", *[f"    • {c}" for c in causes[:5]], ""]
            if impacts:
                lines += ["  IMPACT:", *[f"    • {i}" for i in impacts[:5]], ""]
            if fixes:
                lines += ["  RESOLUTION:", *[f"    • {f}" for f in fixes[:5]], ""]
            if actions:
                lines += ["  ACTION ITEMS:", *[f"    {a}" for a in actions[:8]], ""]

        return "\n".join(lines)

    # ── COMPARISON ────────────────────────────────────────────────────────────

    def _comparison(self, query: str, results: List[dict], cite_map: dict) -> str:
        attr_kws = ["pros","cons","option","chosen","rationale","decision","advantage",
                    "disadvantage","rto","rpo","target","actual","threshold","default",
                    "tier","priority","status","owner","cost","risk","benefit","trade"]
        lines = [f"COMPARISON\n{'─' * 65}",
                 f"Comparing across {min(len(results), self.DOCS_IN_ANSWER)} sources\n"]

        for i, r in enumerate(results[:self.DOCS_IN_ANSWER], 1):
            title = self._title(r)
            meta  = r["metadata"]
            lines += ["─" * 40, f"[{i}] {title}",
                      f"     File: {r['filename']}  ({r.get('file_format','?').upper()})  |  Folder: {r['topic_folder']}"]
            if meta.get("VERSION"):
                lines.append(f"     Version: {meta['VERSION']}  |  Owner: {meta.get('OWNER','—')}")
            lines.append("")
            for sec in r["top_chunks"][:self.CHUNKS_PER_DOC]:
                if sec.get("title"):
                    lines.append(f"  ▸ {sec['title']}")
                for c in self._clean_lines(sec["content"]):
                    if any(kw in c.lower() for kw in attr_kws):
                        lines.append(f"    → {c}")
                    else:
                        lines.append(f"      {c}")
                lines.append("")

        lines += ["─" * 40, "KEY DIFFERENTIATORS", ""]
        diff_kws  = ["chosen","rejected","decided","versus","instead","over",
                     "advantage","disadvantage","rationale","because","benefit","drawback"]
        seen_diff = set()
        for r in results[:self.DOCS_IN_ANSWER]:
            cite = self._cite(r["filename"], cite_map)
            for chunk in r.get("top_chunks", [])[:2]:
                for line in chunk["content"].split("\n"):
                    l = line.strip()
                    if l and any(kw in l.lower() for kw in diff_kws) and l not in seen_diff:
                        seen_diff.add(l)
                        lines.append(f"  {cite} {l}")
        return "\n".join(lines)

    # ── SUMMARY ───────────────────────────────────────────────────────────────

    def _summary(self, query: str, results: List[dict], cite_map: dict) -> str:
        lines = [f"SUMMARY\n{'─' * 65}"]
        for r in results[:self.DOCS_IN_ANSWER]:
            cite  = self._cite(r["filename"], cite_map)
            title = self._title(r)
            meta  = r["metadata"]
            topic = meta.get("TOPIC", r["topic_folder"])
            owner = meta.get("OWNER", "—")
            ver   = meta.get("VERSION", "")
            upd   = meta.get("LAST_UPDATED", "")
            fmt   = r.get("file_format", ".txt").upper()
            lines += ["─" * 40, f"{cite} {title}  [{fmt}]",
                      f"  Topic  : {topic}",
                      f"  Owner  : {owner}" + (f"  |  v{ver}" if ver else "") + (f"  |  Updated: {upd}" if upd else ""),
                      ""]
            for chunk in r["top_chunks"][:self.CHUNKS_PER_DOC]:
                if chunk.get("title"):
                    lines.append(f"  [{chunk['title']}]")
                for l in self._clean_lines(chunk["content"]):
                    lines.append(f"    {l}")
                lines.append("")
        return "\n".join(lines)

    # ── LIST ──────────────────────────────────────────────────────────────────

    def _list(self, query: str, results: List[dict], cite_map: dict) -> str:
        """
        Extracts structured list items matching known prefixes (INC-XXXX, REL-XXX, etc.)
        from the raw document text.

        v7 FIX: fallback behaviour changed for open-ended risk/approval queries.
        Previously: fell back to _summary() if no prefixed items found.
        Now:
          - If the query contains risk/approval vocabulary AND no structured items were
            found → route to _analytical(), which does full content extraction.
          - If structured items were found AND the query is a risk/approval query →
            render both the structured list AND the analytical content.
          - If no risk vocabulary → behave as before (fall back to _summary()).

        This ensures "what risks should I know before approving X" gets the substantive
        risk content (Redis TTL constraints, FRAUD_KAFKA_CONSUMER_GROUP, null rate gate)
        rather than only the INC-/REL- prefixed items that LIST was matching.
        """
        lines    = [f"RESULTS\n{'─' * 65}"]
        item_pat = re.compile(
            r'^\s*(APP:|PORT\s+\d+|PARAM:|ROLE:|TIER\s+\d|JOB:|SEGMENT:|LAYER\s+\d+|'
            r'MILESTONE:|INITIATIVE:|RISK\s+\d|PATTERN\s+\d|TOPIC:|BREACH:|PRIORITY\s+\d|'
            r'VENDOR:|PRODUCT:|INC-\d+|REL-\d+|ADR-\d+|STAGE\s+\d)'
        )

        # ── v7 FIX: risk/approval query detection ───────────────────────────
        RISK_VOCAB = {
            "risk", "risks", "before", "approving", "approve", "approval",
            "gate", "concern", "concerns", "should", "know", "ready",
            "safe", "constraint", "constraints", "issue", "issues",
            "gotcha", "gotchas", "caveat", "caveats",
        }
        query_words   = set(query.lower().split())
        is_risk_query = bool(query_words & RISK_VOCAB)
        # ── end v7 FIX ──────────────────────────────────────────────────────

        by_source: Dict[str, List[str]] = defaultdict(list)
        seen: Set[str] = set()

        for r in results[:6]:
            for line in r["raw_text"].split("\n"):
                s = line.strip()
                if not s or re.match(r"^\s*={3,}", s):
                    continue
                if item_pat.match(s) and s not in seen:
                    seen.add(s)
                    by_source[r["filename"]].append(s)

        if by_source:
            for fname, items in by_source.items():
                r    = next((x for x in results if x["filename"] == fname), None)
                fmt  = r.get("file_format", "").upper() if r else ""
                cite = self._cite(fname, cite_map)
                lines.append(f"\n  {cite} [{fname}  {fmt}]")
                for item in items[:25]:
                    lines.append(f"    • {item}")
            # For risk queries: also append full analytical content beneath the list
            if is_risk_query:
                lines.append(f"\n{'─' * 65}")
                lines.append("ADDITIONAL RISK CONTEXT FROM DOCUMENTS")
                lines.append("─" * 65)
                analytical_body = self._analytical(query, results, cite_map)
                # Strip the "ANALYTICAL BRIEF" header — we already have context
                analytical_lines = analytical_body.split("\n")
                lines.extend(analytical_lines[2:])   # skip first two header lines
            return "\n".join(lines)
        else:
            # No structured items found
            if is_risk_query:
                # Risk query with no prefixed items → full analytical treatment
                return self._analytical(query, results, cite_map)
            else:
                return self._summary(query, results, cite_map)

    # ── TROUBLESHOOT ──────────────────────────────────────────────────────────

    def _troubleshoot(self, query: str, results: List[dict], cite_map: dict) -> str:
        step_kws  = ["check","verify","run","curl","ssh","restart","inspect","test",
                     "confirm","clear","flush","trigger","review","look","ping"]
        cause_kws = ["root cause","common cause","contributing","reason","caused by",
                     "due to","because","failure"]
        cmd_kws   = ["curl","ssh","systemctl","psql","redis-cli","kafka-","kubectl",
                     "python","grep","tail","cat","ls","df","top","htop","ldapsearch"]
        esc_kws   = ["escalat","pagerduty","oncall","@company","page","sev-","notify"]

        lines          = [f"TROUBLESHOOTING GUIDE\n{'─' * 65}"]
        all_escalation = []

        for r in results[:self.DOCS_IN_ANSWER]:
            cite  = self._cite(r["filename"], cite_map)
            title = self._title(r)
            lines += [f"\n{'─' * 40}",
                      f"{cite} {title}  [{r['topic_folder']}]  ({r.get('file_format','').upper()})", ""]

            for chunk in r["top_chunks"][:self.CHUNKS_PER_DOC]:
                causes, steps, commands, esc = [], [], [], []
                for line in chunk["content"].split("\n"):
                    l = line.strip()
                    if not l:
                        continue
                    l_low = l.lower()
                    if any(kw in l_low for kw in cause_kws):
                        causes.append(l)
                    elif any(kw in l_low for kw in cmd_kws):
                        commands.append(l)
                    elif any(kw in l_low for kw in step_kws):
                        steps.append(l)
                    if any(kw in l_low for kw in esc_kws) and l not in all_escalation:
                        esc.append(l)
                        all_escalation.append(l)

                if causes or steps or commands:
                    if chunk.get("title"):
                        lines.append(f"  Section: {chunk['title']}")
                    if causes:
                        lines.append("  Root Cause(s):")
                        for c in causes[:4]: lines.append(f"    • {c}")
                    if steps:
                        lines.append("  Diagnostic Steps:")
                        for s in steps[:15]: lines.append(f"    {s}")
                    if commands:
                        lines.append("  Commands:")
                        for cmd in commands[:8]: lines.append(f"    > {cmd}")
                    lines.append("")

        if all_escalation:
            lines += ["─" * 40, "ESCALATION PATH"]
            for e in list(dict.fromkeys(all_escalation))[:6]:
                lines.append(f"  {e}")

        return "\n".join(lines)

    # ── ANALYTICAL ────────────────────────────────────────────────────────────

    def _analytical(self, query: str, results: List[dict], cite_map: dict) -> str:
        lines = [f"ANALYTICAL BRIEF\n{'─' * 65}", f"Query: {query}", ""]
        lines += ["FINDINGS BY SOURCE", "─" * 40]

        for r in results[:self.DOCS_IN_ANSWER]:
            cite  = self._cite(r["filename"], cite_map)
            title = self._title(r)
            meta  = r["metadata"]
            lines += [f"\n{cite} {title}  [{r.get('file_format','').upper()}]  [{r['topic_folder']}]",
                      f"  Owner: {meta.get('OWNER','—')}  |  Updated: {meta.get('LAST_UPDATED','—')}", ""]
            for chunk in r["top_chunks"][:self.CHUNKS_PER_DOC]:
                if chunk.get("title"):
                    lines.append(f"  [{chunk['title']}]")
                for l in self._clean_lines(chunk["content"]):
                    lines.append(f"    {l}")
                lines.append("")

        # Numeric data points: pull lines with measurable values
        num_pat    = re.compile(r'\d[\d,\.]*\s*(%|ms|GB|TB|TPS|min|hr|month|year|sec|\$|days|K|M)')
        data_points = []
        seen_dp    = set()
        for r in results[:self.DOCS_IN_ANSWER]:
            for line in r["raw_text"].split("\n")[12:]:
                l = line.strip()
                if num_pat.search(l) and l and l not in seen_dp:
                    if not re.match(r"^\s*={3,}", l) and not l.startswith("|"):
                        seen_dp.add(l)
                        data_points.append((l, r["filename"]))

        if data_points:
            lines += ["─" * 40, "KEY METRICS AND DATA POINTS", ""]
            for dp, src in data_points[:20]:
                cite = self._cite(src, cite_map)
                lines.append(f"  • {dp}")
                lines.append(f"      {cite} [{src}]")

        # Patterns and recommendations
        pattern_kws  = ["pattern","lesson","recommendation","risk","trend","action",
                        "priority","proposed","open","in progress","warning","critical"]
        pattern_lines = []
        seen_pl       = set()
        for r in results[:self.DOCS_IN_ANSWER]:
            matches = self._extract_matching(r["raw_text"], pattern_kws, ctx=2)
            for m in matches[:3]:
                if m not in seen_pl:
                    seen_pl.add(m)
                    pattern_lines.append((m, r["filename"]))

        if pattern_lines:
            lines += ["\n─" * 40, "PATTERNS, RISKS AND RECOMMENDATIONS", ""]
            for pl, src in pattern_lines[:10]:
                cite = self._cite(src, cite_map)
                lines.append(f"  {cite} [{src}]")
                for l in pl.split("\n"):
                    lines.append(f"    {l.strip()}")
                lines.append("")

        return "\n".join(lines)

    # ── DEPENDENCY ────────────────────────────────────────────────────────────

    def _dependency(self, query: str, results: List[dict], cite_map: dict) -> str:
        dep_kws  = ["depends on","depend","calls","connects","uses","requires",
                    "upstream","downstream","integrates","host:","port","api:","→","->"]
        app_kws  = ["app:","application","host:","port:","api:","tech stack","endpoint"]
        flow_kws = ["→","->","chain","flow","route","sequence","order"]

        lines = [f"DEPENDENCY MAP\n{'─' * 65}"]
        for r in results[:self.DOCS_IN_ANSWER]:
            cite  = self._cite(r["filename"], cite_map)
            title = self._title(r)
            lines += [f"\n{'─' * 40}", f"{cite} {title}  [{r['topic_folder']}]", ""]

            app_lines, dep_lines, flow_lines = [], [], []
            for chunk in r["top_chunks"][:self.CHUNKS_PER_DOC]:
                for line in chunk["content"].split("\n"):
                    l = line.strip()
                    if not l or re.match(r"^\s*={3,}", l):
                        continue
                    l_low = l.lower()
                    if any(kw in l_low for kw in app_kws):  app_lines.append(l)
                    if any(kw in l_low for kw in dep_kws):  dep_lines.append(l)
                    if any(kw in l_low for kw in flow_kws): flow_lines.append(l)

            if app_lines:
                lines.append("  Application Details:")
                for a in list(dict.fromkeys(app_lines))[:15]: lines.append(f"    {a}")
                lines.append("")
            if dep_lines:
                lines.append("  Dependencies and Connections:")
                for d in list(dict.fromkeys(dep_lines))[:20]: lines.append(f"    {d}")
                lines.append("")
            if flow_lines:
                lines.append("  Data / Request Flow:")
                for f in list(dict.fromkeys(flow_lines))[:8]: lines.append(f"    {f}")
                lines.append("")

        return "\n".join(lines)


# =============================================================================
# PHASE 7 — OUTPUT WRITER
# =============================================================================

class OutputWriter:
    """Persists each query result as a timestamped .txt file and maintains session history."""

    def __init__(self, outputs_dir: str):
        self.dir = Path(outputs_dir)
        self.dir.mkdir(parents=True, exist_ok=True)

    def save(self, result: dict) -> str:
        ts     = datetime.now().strftime("%Y%m%d_%H%M%S")
        slug   = re.sub(r'[^a-z0-9]+', '_', result["query"].lower())[:50].strip("_")
        intent = result.get("intent", "query").lower()
        conf   = result.get("confidence", {}).get("level", "").lower()
        fname  = f"{ts}_{intent}_{conf}_{slug}.txt"
        fpath  = self.dir / fname

        header = (
            f"ENTERPRISE KNOWLEDGE ASSISTANT v7 — QUERY OUTPUT\n"
            f"{'=' * 65}\n"
            f"Query      : {result['query']}\n"
            f"Intent     : {result['intent']}\n"
            f"Confidence : {result.get('confidence',{}).get('level','—')} "
            f"({result.get('confidence',{}).get('score','—')}%)\n"
            f"Sources    : {len(result['sources'])} document(s)\n"
            f"Timestamp  : {result['timestamp']}\n"
            f"{'=' * 65}\n\n"
        )
        with open(fpath, "w", encoding="utf-8") as f:
            f.write(header)
            f.write(result["answer"])
        logger.info(f"  Output: {fpath.name}")
        return str(fpath)


# =============================================================================
# MASTER INTERFACE
# =============================================================================

class EnterpriseKnowledgeAssistant:
    """
    v7 — Four targeted bug fixes + 8 new unstructured knowledge base documents.
    Architecture is unchanged from v6: same 7-phase pipeline, same class hierarchy.

    Startup sequence:
      1. Load manual synonym registry
      2. Ingest all documents (multi-format: .txt/.docx/.xlsx/.pptx)
      3. Auto-synonym discovery across all raw text
      4. Merge discovered synonyms (manual always wins)
      5. Export discovered synonyms to auto_synonyms_discovered.yaml for review
      6. Build BM25 index (chunk-level + document-level)
      7. Ready

    Usage:
        assistant = EnterpriseKnowledgeAssistant()
        assistant.print_answer("What port does the mainframe use?")
        assistant.print_answer("Which system is closest to a capacity breach?")
        assistant.print_answer("What risks should I know before approving fraud engine v3.4.0?")
        assistant.print_answer("Walk me through INC-2024-1201")
    """

    def __init__(self,
                 knowledge_base_dir: str = "./knowledge_base",
                 outputs_dir:        str = "./outputs",
                 synonym_registry:   str = None):

        logger.info("\n" + "=" * 65)
        logger.info("ENTERPRISE KNOWLEDGE ASSISTANT v7")
        logger.info("4 Fixes: Chunk Dedup | Capacity Intent | LIST Risk | Param Noise")
        logger.info("=" * 65)

        reg_path = synonym_registry or str(Path(__file__).parent / "synonym_registry.yaml")

        # Phase 0: Load manual synonym registry
        self.synonyms = SynonymRegistry(reg_path)

        # Phase 1: Ingest all documents
        self.ingestion = DocumentIngestionEngine(knowledge_base_dir, self.synonyms)

        # Auto-synonym discovery: runs on raw text after ingestion, before indexing
        discoverer = AutoSynonymDiscovery(co_occur_threshold=2)
        discovered = discoverer.discover_from_documents(self.ingestion.documents)

        # Export for analyst review
        auto_yaml_path = Path(__file__).parent / "auto_synonyms_discovered.yaml"
        discoverer.export_yaml(auto_yaml_path, discovered)

        # Merge auto-discovered into registry (manual entries always win)
        self.synonyms.merge_auto_discovered(discovered)

        # Phase 2: Build BM25 index
        self.index = BM25IndexEngine(self.ingestion.documents, self.synonyms)

        # Phase 3–7: Retrieval, synthesis, output
        self.retrieval = RetrievalEngine(self.index, self.synonyms)
        self.synthesis = AnswerSynthesisEngine()
        self.writer    = OutputWriter(outputs_dir)
        self.query_history: List[dict] = []

        n      = len(self.ingestion.documents)
        chunks = len(self.index.chunk_corpus)
        fmts   = Counter(d["file_format"] for d in self.ingestion.documents.values())
        quals  = Counter(d["metadata"].get("_quality", "?") for d in self.ingestion.documents.values())

        logger.info(f"\n  ✅ Ready — v7")
        logger.info(f"     Documents   : {n}  |  Chunks: {chunks}")
        logger.info(f"     Formats     : {dict(fmts)}")
        logger.info(f"     Quality     : {dict(quals)}")
        logger.info(f"     Synonyms    : {len(self.synonyms.canonical_to_aliases)} groups "
                    f"({len(self.synonyms.alias_to_canonical)} aliases)")
        logger.info(f"     Auto-YAML   : {auto_yaml_path}")
        logger.info("=" * 65)

    def ask(self, query: str, top_k: int = 6, folder_filter: str = None) -> dict:
        """
        Run a query through the full pipeline: detect intent → retrieve → synthesise → save.
        Returns the full result dict including answer, sources, confidence, and output file path.
        """
        logger.info(f'\n  Query: "{query}"')
        intent  = self.synthesis.intent_detector.detect(query)
        results = self.retrieval.search(query, top_k=top_k,
                                        folder_filter=folder_filter, intent=intent)
        result  = self.synthesis.synthesise(query, results)
        outpath = self.writer.save(result)
        result["output_file"] = outpath
        self.query_history.append({
            "query":       query,
            "intent":      result.get("intent"),
            "confidence":  result.get("confidence", {}).get("level"),
            "timestamp":   result["timestamp"],
            "status":      result["status"],
            "top_source":  results[0]["filename"] if results else None,
            "output_file": outpath,
        })
        return result

    def print_answer(self, query: str, top_k: int = 6, folder_filter: str = None):
        """Run a query and print the formatted answer to stdout."""
        result = self.ask(query, top_k=top_k, folder_filter=folder_filter)
        print(result["answer"])
        print(f"\n  → Saved: {result['output_file']}")
        return result

    def get_document_list(self) -> List[dict]:
        """Return a summary of all ingested documents."""
        return [
            {"doc_id": d["doc_id"], "filename": d["filename"],
             "file_format": d["file_format"], "topic_folder": d["topic_folder"],
             "char_count": d["char_count"],
             "quality": d["metadata"].get("_quality", "?"),
             "metadata": d["metadata"]}
            for d in self.ingestion.documents.values()
        ]

    def get_query_history(self) -> List[dict]:
        """Return the session query history list."""
        return self.query_history

    def export_query_history(self, filepath: str = "query_history.json"):
        """Write session query history to a JSON file."""
        with open(filepath, "w") as f:
            json.dump(self.query_history, f, indent=2, default=str)
        logger.info(f"  History → {filepath}")

    def get_auto_synonyms(self) -> dict:
        """Return only the auto-discovered synonym groups (excludes manual ones)."""
        manual = self.synonyms._manual_canonicals
        return {
            k: v for k, v in self.synonyms.canonical_to_aliases.items()
            if k not in manual
        }
